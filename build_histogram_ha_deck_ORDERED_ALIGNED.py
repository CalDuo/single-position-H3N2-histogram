#!/usr/bin/env python3
"""
Build a PowerPoint deck with one slide per position:
  - Histogram on the left: width 7.5 inches, left edge flush to slide
  - HA-head on the right: width 3.0 inches, right edge flush to slide
  - Both vertically centered
  - Slide order: descending numeric position
  - Treat negative IDs as the same as positive (e.g., -117 == 117)

Inputs:
  --hist_zip : ZIP containing histogram SVG files (any folder structure)
  --ha_zip   : ZIP containing HA-head PNG files (any folder structure)
  --out      : output PPTX file

Notes:
  - PowerPoint embedding requires raster; histogram SVGs are converted to PNGs via cairosvg.
  - HA-head selection: if multiple PNGs per position exist, choose by priority:
      purple > red > blue > green > otherwise first found.
"""

from __future__ import annotations

import argparse
import re
import zipfile
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import cairosvg
from pptx import Presentation
from pptx.util import Inches


def extract_pos(name: str) -> Optional[int]:
    m = re.search(r"-?\d{1,3}", name)
    if not m:
        return None
    return abs(int(m.group(0)))


def unzip(zip_path: Path, out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(zip_path, "r") as z:
        z.extractall(out_dir)


def convert_svgs_to_pngs(svg_root: Path, png_dir: Path, width_px: int = 1800) -> Dict[int, Path]:
    png_dir.mkdir(parents=True, exist_ok=True)
    out: Dict[int, Path] = {}
    for pth in svg_root.rglob("*.svg"):
        pos = extract_pos(pth.name)
        if pos is None:
            continue
        out_png = png_dir / f"{pos}.png"
        cairosvg.svg2png(url=str(pth), write_to=str(out_png), output_width=width_px)
        out[pos] = out_png
    return out


def choose_ha_pngs(ha_root: Path) -> Dict[int, Path]:
    priority = ["purple", "red", "blue", "green"]
    candidates: Dict[int, List[Tuple[int, Path]]] = {}
    for pth in ha_root.rglob("*.png"):
        pos = extract_pos(pth.name)
        if pos is None:
            continue
        name = pth.name.lower()
        score = 100
        for i, key in enumerate(priority):
            if key in name:
                score = i
                break
        candidates.setdefault(pos, []).append((score, pth))

    chosen: Dict[int, Path] = {}
    for pos, lst in candidates.items():
        chosen[pos] = sorted(lst, key=lambda x: x[0])[0][1]
    return chosen


def build_deck(hist_pngs: Dict[int, Path], ha_pngs: Dict[int, Path], out_pptx: Path) -> None:
    positions = sorted(set(hist_pngs.keys()) & set(ha_pngs.keys()), reverse=True)

    prs = Presentation()
    blank = prs.slide_layouts[6]

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Approved layout rules
    hist_width = Inches(7.5)
    hist_left = Inches(0.0)

    ha_width = Inches(3.0)
    ha_left = slide_width - ha_width  # flush right

    for pos in positions:
        slide = prs.slides.add_slide(blank)

        pic_hist = slide.shapes.add_picture(str(hist_pngs[pos]), hist_left, Inches(0.4), width=hist_width)
        pic_hist.top = int((slide_height - pic_hist.height) / 2)

        pic_ha = slide.shapes.add_picture(str(ha_pngs[pos]), ha_left, Inches(0.4), width=ha_width)
        pic_ha.top = int((slide_height - pic_ha.height) / 2)

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--hist_zip", required=True, type=str, help="ZIP containing histogram SVGs")
    ap.add_argument("--ha_zip", required=True, type=str, help="ZIP containing HA-head PNGs")
    ap.add_argument("--out", required=True, type=str, help="Output PPTX path")
    ap.add_argument("--tmp_dir", default=".ppt_build", type=str, help="Temporary working dir")
    ap.add_argument("--svg_width_px", default=1800, type=int, help="Converted histogram PNG width")
    args = ap.parse_args()

    tmp = Path(args.tmp_dir)
    hist_dir = tmp / "hist"
    ha_dir = tmp / "ha"
    hist_png_dir = tmp / "hist_png"

    unzip(Path(args.hist_zip), hist_dir)
    unzip(Path(args.ha_zip), ha_dir)

    hist_pngs = convert_svgs_to_pngs(hist_dir, hist_png_dir, width_px=args.svg_width_px)
    ha_pngs = choose_ha_pngs(ha_dir)

    build_deck(hist_pngs, ha_pngs, Path(args.out))


if __name__ == "__main__":
    main()
